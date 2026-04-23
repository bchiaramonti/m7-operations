#!/usr/bin/env python3
"""
build_pptx.py — Render 8-slide status presentation (16:9) from canonical data.

Mapping canvas Paper `status-report`:
  01 Cover · 02 Agenda · 03 Roadmap Overview · 04 Roadmap Detail
  05 Section Divider · 06 Executive Status · 07 Risks · 08 Closing

Inputs:
  --data <json-path>       Canonical dict from collect_data.py
  --assets-dir <path>      Directory with m7-logo-*.png
  --out-dir <path>         Where to write status-presentation.pptx

Pure python-pptx — no MCP or LLM calls.
"""

from __future__ import annotations

import argparse
import json
import sys
import tempfile
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("✗ python-pptx ausente. Instale: pip install python-pptx", file=sys.stderr)
    sys.exit(2)

# Import sibling render helpers for HTML screenshot embedding (slide Roadmap)
sys.path.insert(0, str(Path(__file__).parent))
try:
    from render_html_section import (
        render as render_html_screenshot,
        build_roadmap_overlay_js,
        PRESETS as RENDER_PRESETS,
    )
except ImportError:
    render_html_screenshot = None
    build_roadmap_overlay_js = None
    RENDER_PRESETS = {}


# ---- Tokens (mirrored from references/design-tokens.md) ----

COL_PRIMARY        = RGBColor(0x42, 0x41, 0x35)
COL_BG_LIGHT       = RGBColor(0xFF, 0xFD, 0xEF)
COL_TEXT_MUTED     = RGBColor(0x4F, 0x4E, 0x3C)
COL_TEXT_SUBTLE    = RGBColor(0x79, 0x75, 0x5C)
COL_TEXT_CAPTION   = RGBColor(0xAE, 0xAD, 0xA8)
COL_ACCENT_LIME    = RGBColor(0xEE, 0xF7, 0x7C)
COL_STATUS_OK      = RGBColor(0x00, 0xB0, 0x50)
COL_STATUS_PROG    = RGBColor(0x3B, 0x82, 0xF6)
COL_STATUS_WARN    = RGBColor(0xF5, 0x9E, 0x0B)
COL_STATUS_CRIT    = RGBColor(0xE4, 0x69, 0x62)
COL_STATUS_NEUTRAL = RGBColor(0xD0, 0xD0, 0xCC)
COL_TAG_BG_CRIT    = RGBColor(0xFD, 0xED, 0xED)
COL_TAG_BG_WARN    = RGBColor(0xFD, 0xF3, 0xE0)
COL_TAG_BG_NEUTRAL = RGBColor(0xEF, 0xEF, 0xEC)
COL_TABLE_HDR_BG   = RGBColor(0xF1, 0xEF, 0xDE)
COL_ROW_BORDER     = RGBColor(0xE5, 0xE2, 0xCE)

# ---- Slide layout constants (documented values from references/design-tokens.md) ----
# Pixel-to-EMU conversion: 1 px = 9525 EMU (slide 13.333in x 7.5in @ 96dpi = 1280x720 px)
PX = 9525

# Canvas dimensions (Paper artboards are 1280x720; see references/presentation-structure.md)
SLIDE_W_PX = 1280
SLIDE_H_PX = 720

# Paddings (vertical / horizontal) — see design-tokens.md §Espaçamento
PAD_SLIDE_X = 80        # default horizontal padding for content slides
PAD_SLIDE_Y_TOP = 56    # default top padding
PAD_SLIDE_Y_BOT = 56    # default bottom padding
PAD_EXEC_X = 40         # slide 06 Executive uses tighter padding (denser)
PAD_EXEC_Y = 36

# Content column width (after left+right padding): 1280 - 80*2 = 1120
CONTENT_W_PX = SLIDE_W_PX - PAD_SLIDE_X * 2

# Milestone/timeline element sizes — see design-tokens.md §Elementos estruturais
DIAMOND_SIZE_ACTIVE = 14    # filled diamond (done, in_progress, overdue)
DIAMOND_SIZE_NEUTRAL = 10   # outlined diamond (not_started)

# Risk card dimensions — see presentation-structure.md §Slide 07
RISK_CARD_H = 120
RISK_CARD_GAP = 14
RISK_ACCENT_BAR_W = 6
RISK_TAG_W = 150

# Logo display sizes (source PNG is 196x96; aspect preserved)
LOGO_SIZE_COVER = 60        # Cover, Agenda, Section Divider, Closing (dark BGs)
LOGO_SIZE_CONTENT = 60      # Roadmap, Risks (light BGs)
LOGO_SIZE_EXEC = 44         # Executive Status (tighter padding)


def emu(px_val: float) -> int:
    return int(px_val * PX)


# ---- Helpers ----

def fill_solid(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_bg_rect(slide, w_emu, h_emu, rgb: RGBColor):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w_emu, h_emu)
    fill_solid(rect, rgb)


def add_bg_image(slide, w_emu, h_emu, image_path: Path) -> bool:
    """Adds an image as full-bleed slide background. Returns False if path missing."""
    if not image_path.exists():
        return False
    slide.shapes.add_picture(str(image_path), 0, 0, width=w_emu, height=h_emu)
    return True


def add_bg_overlay(slide, w_emu, h_emu, rgb: RGBColor, transparency: float = 0.0):
    """Adds a rectangle overlay (z-ordered above background image) for darkening/tinting.

    transparency is 0.0 (opaque) to 1.0 (fully transparent). Default opaque.
    """
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w_emu, h_emu)
    rect.fill.solid()
    rect.fill.fore_color.rgb = rgb
    rect.line.fill.background()
    if transparency > 0.0:
        from pptx.oxml.ns import qn
        # python-pptx doesn't expose transparency directly; inject via raw XML
        solid_fill = rect.fill._xPr.find(qn("a:solidFill"))
        if solid_fill is not None:
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is not None:
                from lxml import etree
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", str(int((1.0 - transparency) * 100000)))
    return rect


def add_text(
    slide, x_px, y_px, w_px, h_px, text, *,
    size=14, bold=False, color=COL_PRIMARY, tracking=0.0,
    align="left", anchor="top", italic=False, font="Arial"
):
    box = slide.shapes.add_textbox(emu(x_px), emu(y_px), emu(w_px), emu(h_px))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    # First paragraph
    p = tf.paragraphs[0]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    for i, line in enumerate(str(text).split("\n")):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if tracking:
            # spc is in 1/100 pt; 1em ~= size pt, so spc = tracking * size * 100
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(tracking * size * 100)))
    return box


def add_rect(slide, x_px, y_px, w_px, h_px, rgb: RGBColor, line_rgb=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x_px), emu(y_px), emu(w_px), emu(h_px))
    fill_solid(r, rgb)
    if line_rgb is not None:
        r.line.color.rgb = line_rgb
        r.line.width = Emu(0)
    return r


def add_diamond(slide, x_px, y_px, size_px, rgb: RGBColor, outlined=False):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, emu(x_px), emu(y_px), emu(size_px), emu(size_px))
    if outlined:
        d.fill.background()
        d.line.color.rgb = rgb
        d.line.width = Pt(1.5)
    else:
        fill_solid(d, rgb)
    return d


def add_logo(slide, x_px, y_px, size_px, assets_dir: Path, variant: str):
    path = assets_dir / f"m7-logo-{variant}.png"
    if path.exists():
        # Logo original is 196x96; maintain aspect ratio, target width~44-56px
        slide.shapes.add_picture(str(path), emu(x_px), emu(y_px), width=emu(size_px))
    else:
        color = COL_BG_LIGHT if variant == "offwhite" else COL_PRIMARY
        add_text(slide, x_px, y_px, size_px, size_px * 0.45, "M7",
                 size=22, bold=True, italic=True, color=color, align="right")


# ---- Severity helpers ----

def risk_severity_label(risk) -> tuple[str, RGBColor, RGBColor]:
    """Returns (label, bg_color, bar_color) for a risk card.

    Uses the authoritative severity_class from riscos.html when present (crit/high/med/low),
    falls back to probability × impact derivation for legacy data.
    """
    sc = risk.get("severity_class") or ""
    prob = (risk.get("probability") or "").lower()
    impact = (risk.get("impact") or "").lower()

    if sc == "crit":
        return f"{prob.upper()} · {_impact_label(impact)}", COL_TAG_BG_CRIT, COL_STATUS_CRIT
    if sc == "high":
        return f"{prob.upper()} · {_impact_label(impact)}", COL_TAG_BG_WARN, COL_STATUS_WARN
    if sc == "med":
        return f"{prob.upper()} · {_impact_label(impact)}", COL_TAG_BG_NEUTRAL, COL_STATUS_WARN
    if sc == "low":
        return f"{prob.upper()} · {_impact_label(impact)}", COL_TAG_BG_NEUTRAL, COL_STATUS_NEUTRAL

    # Legacy fallback
    if prob == "alta" and impact in ("critico", "alto"):
        return "ALTA · CRÍTICO", COL_TAG_BG_CRIT, COL_STATUS_CRIT
    if prob in ("alta", "media") or impact in ("critico", "alto"):
        return f"{prob.upper()} · {_impact_label(impact)}", COL_TAG_BG_WARN, COL_STATUS_WARN
    return f"{prob.upper()} · {_impact_label(impact)}", COL_TAG_BG_NEUTRAL, COL_STATUS_NEUTRAL


def _impact_label(impact: str) -> str:
    return {
        "critico": "CRÍTICO", "alto": "ALTO",
        "medio": "MÉDIO", "baixo": "BAIXO",
    }.get((impact or "").lower(), impact.upper())


# ---- Slides ----

def slide_01_cover(prs, data, assets_dir):
    """Cover fiel ao artboard Paper `01 — Cover`:

    Layer stack (z-order bottom → top):
      1. hero photo (m7-hero-dark.png, full-bleed)
      2. dark tint overlay (~35% black) to ensure text readability on photo
      3. text/logo elements
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background — prefer hero photo; fallback to solid caqui if asset missing
    hero = assets_dir / "m7-hero-dark.png"
    if not add_bg_image(slide, prs.slide_width, prs.slide_height, hero):
        add_bg_rect(slide, prs.slide_width, prs.slide_height, COL_PRIMARY)
    else:
        # Tint overlay matching Paper gradient (55–75% black in oklab ≈ #424135 @ 35% transparency)
        add_bg_overlay(slide, prs.slide_width, prs.slide_height, COL_PRIMARY, transparency=0.35)

    # Eyebrow "STATUS REPORT" — lime accent, tracked, top-left
    add_text(slide, 80, 60, 400, 20, "STATUS REPORT",
             size=13, bold=True, color=COL_ACCENT_LIME, tracking=0.2)

    # M7 logo top-right (offwhite variant for dark BG)
    add_logo(slide, 1180, 48, 60, assets_dir, "offwhite")

    # Title block — project name + "Status Report" stacked, positioned per Paper artboard
    project_name = data.get("project", {}).get("name", "Projeto")
    add_text(slide, 80, 290, 1120, 90, project_name,
             size=64, bold=False, color=COL_BG_LIGHT)
    add_text(slide, 80, 370, 1120, 90, "Status Report",
             size=64, bold=False, color=COL_BG_LIGHT)

    # Subtitle (period label) in lime accent
    add_text(slide, 80, 470, 960, 30,
             data.get("project", {}).get("period_label", ""),
             size=20, color=COL_ACCENT_LIME)

    # Footer: subtle divider + M7 tagline left + slide number right
    add_rect(slide, 80, 640, 1120, 1, COL_TEXT_CAPTION)
    add_text(slide, 80, 660, 800, 20,
             data.get("project", {}).get("footer_label", ""),
             size=12, color=COL_BG_LIGHT)
    add_text(slide, 1160, 660, 80, 20, "01",
             size=12, color=COL_BG_LIGHT, align="right")


def slide_02_agenda(prs, data, assets_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, prs.slide_width, prs.slide_height, COL_PRIMARY)
    add_text(slide, 80, 64, 400, 20, "STATUS REPORT",
             size=13, bold=True, color=COL_ACCENT_LIME, tracking=0.2)
    add_logo(slide, 1180, 50, 60, assets_dir, "offwhite")
    # Title
    add_text(slide, 80, 140, 800, 60, "Agenda",
             size=48, color=COL_BG_LIGHT)
    # Divider
    add_rect(slide, 80, 220, 1120, 1, COL_BG_LIGHT)
    # Items — deck de 7 slides (Cover · Agenda · Visão Geral · Roadmap · Mapa Status · Riscos · Closing)
    items = [
        ("01", "Visão Geral do Roadmap", "Slide 03"),
        ("02", "Roadmap Completo", "Slide 04"),
        ("03", "Mapa de Status Executivo", "Slide 05"),
        ("04", "Riscos Ativos", "Slide 06"),
        ("05", "Próximos Passos", "Slide 07"),
    ]
    y = 246
    for num, title, rng in items:
        add_text(slide, 80, y, 60, 20, num,
                 size=13, bold=True, color=COL_ACCENT_LIME, tracking=0.2)
        add_text(slide, 160, y, 800, 28, title,
                 size=22, color=COL_BG_LIGHT)
        add_text(slide, 980, y + 4, 220, 20, rng,
                 size=10, color=COL_TEXT_CAPTION, align="right")
        y += 40
        add_rect(slide, 80, y + 16, 1120, 1, COL_BG_LIGHT)
        y += 30


def _render_roadmap_screenshot(
    roadmap_html: Path | None,
    preset: str,
    tmp_dir: Path,
    overlays: dict | None = None,
) -> Path | None:
    """Renders a section of the roadmap-marcos.html as a PNG via playwright.

    When `overlays` is provided (with keys today_pct, today_label, bar_status_by_title),
    the preset `roadmap-full-with-status-overlays` gets extra_js composed from that
    data to dynamically color bars by status and draw the HOJE vertical line.

    Returns the output path or None on failure.
    """
    if render_html_screenshot is None or roadmap_html is None or not roadmap_html.exists():
        return None
    try:
        out = tmp_dir / f"roadmap-{preset}.png"
        p = RENDER_PRESETS.get(preset)
        if not p:
            return None
        extra_js = ""
        if overlays and build_roadmap_overlay_js and preset.endswith("with-status-overlays"):
            extra_js = build_roadmap_overlay_js(
                today_pct=overlays.get("today_pct"),
                today_label=overlays.get("today_label", ""),
                bar_status_map=overlays.get("bar_status_by_title", {}),
            )
        render_html_screenshot(
            input_path=roadmap_html,
            output_path=out,
            selector=p["selector"],
            viewport=p["viewport"],
            device_scale=p["device_scale"],
            inject_css=p["inject_css"],
            extra_js=extra_js,
        )
        return out if out.exists() else None
    except Exception as e:
        print(f"⚠ Screenshot roadmap ({preset}) falhou: {e}", file=sys.stderr)
        return None


def slide_03_visao_geral_roadmap(prs, data, assets_dir, ctx):
    """Slide 3 — Visão Geral do Roadmap (Paper artboard `03 — Visão Geral`).

    Renders a processos × fases matrix:
      - Rows: processos from matrix-structure.json (11 playbooks for this project)
      - Columns: fases do trabalho (Mapa N2, Mapa N3, DEIP, Políticas, Playbook, Implementação)
      - Cells: status-colored squares per task in the (processo, fase) intersection

    Structure is inferred once (persisted to <proj>/4-status-report/matrix-structure.json)
    and can be edited by the user. When inference fails, the JSON has
    source="pending_user_input" and the slide renders a helpful placeholder.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, prs.slide_width, prs.slide_height, COL_BG_LIGHT)

    structure = data.get("roadmap_structure") or {}
    processos = structure.get("processos", [])
    fases = structure.get("fases", [])
    matrix = structure.get("matrix", [])
    meta = structure.get("meta", {}) or {}
    source = structure.get("source", "inferred")

    # Header
    add_text(slide, 80, 40, 400, 16, "03 · VISÃO GERAL",
             size=12, bold=True, color=COL_PRIMARY, tracking=0.2)
    add_logo(slide, 1180, 36, 60, assets_dir, "dark")

    done = meta.get("done_cells", 0)
    active = meta.get("active_cells", 0)
    total = meta.get("total_cells", 0)
    if total > 0:
        pct = round(100 * done / total)
        hero = f"{done + active} de {total} entregas em execução ({pct}% concluídas)"
    else:
        hero = "Visão Geral do Roadmap"
    add_text(slide, 80, 58, 1100, 40, hero, size=26, color=COL_PRIMARY)
    add_text(slide, 80, 100, 1120, 20,
             "Status por processo × fase do trabalho · fonte: Cronograma.xlsx LIVE.",
             size=12, color=COL_TEXT_SUBTLE)

    # Handle pending/empty structure — inform the user
    if source == "pending_user_input" or not processos or not fases:
        add_text(slide, 80, 300, 1120, 30,
                 "Estrutura da matriz ainda não definida.",
                 size=18, bold=True, color=COL_STATUS_WARN, align="center")
        add_text(slide, 80, 340, 1120, 30,
                 "Edite <proj>/4-status-report/matrix-structure.json ou rode o fluxo interativo da skill.",
                 size=12, color=COL_TEXT_SUBTLE, align="center")
        return

    # Matrix layout — adapts column width to number of fases
    n_fases = len(fases)
    table_x = 80
    table_y = 160
    proc_col_w = 260
    # Available width for fase columns:
    avail_w = 1120 - proc_col_w
    fase_col_w = avail_w // max(n_fases, 1)
    header_h = 38
    # Row height adapts to fit up to 12 rows within ~440px of vertical space
    n_rows = len(processos)
    max_rows_space = 440
    row_h = min(40, max_rows_space // max(n_rows, 1))
    cell_square_size = 18

    # Header row — PROCESSO column + one column per fase
    add_rect(slide, table_x, table_y, proc_col_w + n_fases * fase_col_w, header_h,
             COL_TABLE_HDR_BG)
    add_text(slide, table_x + 16, table_y + 11, proc_col_w - 16, 18,
             "PROCESSO", size=10, bold=True, color=COL_TEXT_MUTED, tracking=0.14)
    for i, fase in enumerate(fases):
        fx = table_x + proc_col_w + i * fase_col_w
        add_text(slide, fx, table_y + 11, fase_col_w, 18,
                 fase.get("label", "").upper(), size=10, bold=True,
                 color=COL_TEXT_MUTED, tracking=0.14, align="center")

    # Data rows — one per processo
    status_color = {
        "done": COL_STATUS_OK,
        "active": COL_STATUS_PROG,
        "overdue": COL_STATUS_CRIT,
        "not_started": COL_STATUS_NEUTRAL,
    }
    for row_idx, processo in enumerate(processos):
        ry = table_y + header_h + row_idx * row_h
        # Row bottom border (subtle)
        add_rect(slide, table_x, ry + row_h - 1,
                 proc_col_w + n_fases * fase_col_w, 1, COL_ROW_BORDER)
        # Processo label: small WBS prefix + name
        wbs = processo.get("wbs_prefix", "")
        name = processo.get("label", "")
        if wbs:
            add_text(slide, table_x + 16, ry + row_h // 2 - 10, 52, 20,
                     wbs, size=10, bold=True, color=COL_TEXT_SUBTLE)
            add_text(slide, table_x + 72, ry + row_h // 2 - 10, proc_col_w - 72, 20,
                     truncate(name, 32), size=12, color=COL_TEXT_MUTED)
        else:
            add_text(slide, table_x + 16, ry + row_h // 2 - 10, proc_col_w - 16, 20,
                     truncate(name, 36), size=12, color=COL_TEXT_MUTED)
        # Cells
        row_cells = matrix[row_idx] if row_idx < len(matrix) else []
        for i in range(n_fases):
            cell = row_cells[i] if i < len(row_cells) else {"status": "missing"}
            cx_center = table_x + proc_col_w + i * fase_col_w + fase_col_w // 2
            cy_center = ry + row_h // 2
            cell_status = cell.get("status", "missing")
            if cell_status == "missing":
                add_text(slide, cx_center - 30, cy_center - 9, 60, 18,
                         "—", size=16, color=COL_TEXT_CAPTION, align="center")
            else:
                color = status_color.get(cell_status, COL_STATUS_NEUTRAL)
                add_rect(slide,
                         cx_center - cell_square_size // 2,
                         cy_center - cell_square_size // 2,
                         cell_square_size, cell_square_size, color)

    # Legend + footer
    legend_y = 650
    legend_items = [
        ("Concluída", COL_STATUS_OK),
        ("Em andamento", COL_STATUS_PROG),
        ("Atrasada", COL_STATUS_CRIT),
        ("Não iniciada", COL_STATUS_NEUTRAL),
    ]
    lx = 80
    for label, color in legend_items:
        add_rect(slide, lx, legend_y + 4, 12, 12, color)
        add_text(slide, lx + 18, legend_y + 2, 110, 14,
                 label, size=9, color=COL_TEXT_MUTED)
        lx += 140
    add_text(slide, lx, legend_y + 2, 140, 14,
             "—  Sem task mapeada", size=9, color=COL_TEXT_CAPTION)
    add_text(slide, 80, 680, 1120, 14,
             f"Fonte: Cronograma.xlsx LIVE + matrix-structure.json · atualizado em {data.get('report_date','')}",
             size=9, color=COL_TEXT_CAPTION)


def slide_04_roadmap(prs, data, assets_dir, ctx):
    """Slide 4 — Roadmap Completo: full swim-lane screenshot from roadmap-marcos.html
    with status-based bar coloring and the HOJE vertical reference line.

    Visual source of truth is the same HTML approved when building the plan.
    Bar colors reflect task execution status aggregated from Cronograma.xlsx LIVE
    (see collect_data.py::aggregate_bar_status). The HOJE line is computed by
    interpolating between M0 and M7 anchor milestones.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, prs.slide_width, prs.slide_height, COL_BG_LIGHT)

    # Header
    add_text(slide, 80, 40, 400, 16, "04 · ROADMAP",
             size=12, bold=True, color=COL_PRIMARY, tracking=0.2)
    add_logo(slide, 1180, 36, 60, assets_dir, "dark")
    add_text(slide, 80, 58, 900, 36, "Roadmap Completo",
             size=28, color=COL_PRIMARY)
    add_text(slide, 80, 94, 900, 20,
             "Fases · frentes · marcos · status de execução por flotilha (HOJE em destaque).",
             size=12, color=COL_TEXT_SUBTLE)

    # Screenshot of full swim-lane roadmap with status overlays (bar coloring + HOJE line)
    roadmap_png = _render_roadmap_screenshot(
        ctx.get("roadmap_html"),
        "roadmap-full-with-status-overlays",
        ctx["tmp_dir"],
        overlays=data.get("roadmap_overlays"),
    )
    if roadmap_png:
        # Embed image filling most of the slide (leave room for header & footer)
        slide.shapes.add_picture(
            str(roadmap_png),
            emu(20), emu(130),
            width=emu(1240), height=emu(530),
        )
    else:
        add_text(slide, 80, 360, 1120, 40,
                 "— Screenshot do roadmap indisponível (playwright/chromium ausente) —",
                 size=14, color=COL_TEXT_CAPTION, align="center")

    # Footer: legend + source
    add_text(slide, 80, 670, 700, 14,
             "■ Concluída   ■ Em andamento   ■ Atrasada   ■ Futura   │ HOJE",
             size=9, color=COL_TEXT_MUTED)
    add_text(slide, 80, 688, 1120, 14,
             f"Fonte: 1-planning/artefatos/roadmap-marcos.html · atualizado em {data.get('report_date','')}",
             size=9, color=COL_TEXT_CAPTION)


def slide_05_mapa_status_executivo(prs, data, assets_dir, ctx):
    """Slide 5 — Mapa de Status Executivo (Paper artboard `06 — Executive Status`).

    Layout replicando o Paper:
      - Header: eyebrow "04 · STATUS EXECUTIVO" + título evocativo + hero sentence (X/Y tarefas)
      - Zone 1: Cronograma Macro com marcos M0-M7 + linha HOJE vertical
      - Zones 2-3: 2 colunas (STATUS EXECUTIVO + PRÓXIMAS ATIVIDADES)
      - Zone 4: PONTOS DE ATENÇÃO
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, prs.slide_width, prs.slide_height, COL_BG_LIGHT)

    # Header — eyebrow numerado + título explícito + hero sentence
    add_text(slide, 40, 28, 800, 14, "05 · STATUS EXECUTIVO",
             size=11, bold=True, color=COL_PRIMARY, tracking=0.2)
    add_text(slide, 40, 44, 900, 32, "Mapa de Status Executivo",
             size=24, color=COL_PRIMARY)
    add_text(slide, 40, 78, 900, 24,
             data.get("status", {}).get("hero_sentence", ""),
             size=14, color=COL_TEXT_SUBTLE)
    add_logo(slide, 1200, 36, 44, assets_dir, "dark")

    # Zone 1: Cronograma Macro — M0-M7 gates from roadmap-marcos.html (canonical).
    # Each milestone rendered as diamond with date below. HOJE line overlaid.
    zone1_y = 120
    add_text(slide, 40, zone1_y, 600, 14, "CRONOGRAMA MACRO",
             size=10, bold=True, color=COL_TEXT_MUTED, tracking=0.14)
    mm = data.get("macro_milestones", [])
    if mm:
        timeline_y = zone1_y + 30
        strip_x = 80
        strip_w = 1120
        # Bounding box
        add_rect(slide, 40, zone1_y + 16, 1200, 100, COL_BG_LIGHT, line_rgb=COL_ROW_BORDER)

        n = len(mm)
        spacing = strip_w // max(n - 1, 1) if n > 1 else 0

        # HOJE vertical line — position by interpolating report_date between first and last milestone dates.
        report_date = ctx.get("report_date")
        hoje_x = None
        dated = [m for m in mm if m.get("date")]
        if report_date and len(dated) >= 2:
            try:
                first_dt = datetime.strptime(dated[0]["date"], "%Y-%m-%d")
                last_dt = datetime.strptime(dated[-1]["date"], "%Y-%m-%d")
                total_span = (last_dt - first_dt).days
                offset = (report_date - first_dt).days
                if 0 < total_span:
                    ratio = max(0.0, min(1.0, offset / total_span))
                    hoje_x = strip_x + int(strip_w * ratio)
            except (ValueError, TypeError):
                pass

        for i, m in enumerate(mm):
            cx = strip_x + i * spacing
            status = m.get("status", "not_started")
            color = {
                "done": COL_STATUS_OK,
                "in_progress": COL_STATUS_PROG,
                "overdue": COL_STATUS_CRIT,
                "not_started": COL_STATUS_NEUTRAL,
            }.get(status, COL_STATUS_NEUTRAL)
            outlined = status == "not_started"
            ds = 14 if not outlined else 10
            add_diamond(slide, cx - ds // 2, timeline_y + 14, ds, color, outlined)
            # Label (M0 KICKOFF) + date (14/abr) stacked below
            add_text(slide, cx - 60, timeline_y + 32, 120, 14,
                     m.get("label", ""), size=9, bold=True, color=COL_TEXT_MUTED, align="center", tracking=0.1)
            if m.get("date_short"):
                add_text(slide, cx - 60, timeline_y + 46, 120, 14,
                         m["date_short"], size=9, color=COL_TEXT_CAPTION, align="center")
            if i < n - 1:
                add_rect(slide, cx + ds // 2, timeline_y + 20, spacing - ds, 2,
                         color if status in ("done", "in_progress") else COL_STATUS_NEUTRAL)

        # HOJE vertical line (drawn on top of timeline)
        if hoje_x is not None:
            add_rect(slide, hoje_x - 1, timeline_y - 4, 2, 44, COL_STATUS_CRIT)
            add_text(slide, hoje_x - 30, timeline_y - 18, 60, 12,
                     "HOJE", size=8, bold=True, color=COL_STATUS_CRIT, align="center", tracking=0.2)

    # Zone 2-3: Two columns
    zone2_y = 260
    col_w = 586
    # Label columns
    add_text(slide, 40, zone2_y, col_w, 14, "STATUS EXECUTIVO",
             size=10, bold=True, color=COL_TEXT_MUTED, tracking=0.14)
    add_text(slide, 40 + col_w + 20, zone2_y, col_w, 14, "PRÓXIMAS ATIVIDADES",
             size=10, bold=True, color=COL_TEXT_MUTED, tracking=0.14)

    # Box Left (Highlights)
    add_rect(slide, 40, zone2_y + 20, col_w, 240, COL_BG_LIGHT, line_rgb=COL_ROW_BORDER)
    highlights = data.get("highlights", [])[:5]
    y = zone2_y + 36
    for h in highlights:
        add_rect(slide, 56, y + 6, 4, 4, COL_PRIMARY)
        add_text(slide, 68, y, col_w - 36, 20, h,
                 size=11, color=COL_TEXT_MUTED)
        y += 24

    # Box Right (Next steps)
    add_rect(slide, 40 + col_w + 20, zone2_y + 20, col_w, 240, COL_BG_LIGHT, line_rgb=COL_ROW_BORDER)
    next_steps = data.get("next_steps", [])[:5]
    y = zone2_y + 36
    for n in next_steps:
        add_rect(slide, 40 + col_w + 20 + 16, y + 6, 4, 4, COL_PRIMARY)
        add_text(slide, 40 + col_w + 20 + 28, y, col_w - 36, 20,
                 f"{n.get('action','')} — até {n.get('deadline','')}",
                 size=11, color=COL_TEXT_MUTED)
        y += 24

    # Zone 4: Attentions
    zone4_y = 540
    add_text(slide, 40, zone4_y, 800, 14, "PONTOS DE ATENÇÃO",
             size=10, bold=True, color=COL_STATUS_WARN, tracking=0.14)
    add_rect(slide, 40, zone4_y + 20, 1200, 90, COL_BG_LIGHT, line_rgb=COL_ROW_BORDER)
    y = zone4_y + 32
    for a in data.get("attentions", [])[:3]:
        sev_color = {
            "critical": COL_STATUS_CRIT,
            "warning": COL_STATUS_WARN,
            "neutral": COL_STATUS_NEUTRAL,
        }.get(a.get("severity"), COL_STATUS_NEUTRAL)
        add_rect(slide, 56, y + 6, 4, 4, sev_color)
        add_text(slide, 68, y, 1130, 20, a.get("text", ""),
                 size=11, color=COL_TEXT_MUTED)
        y += 24

    # Footer
    add_text(slide, 40, 680, 600, 14,
             f"Atualizado em {data.get('report_date', '')} · Fonte: Cronograma.xlsx LIVE",
             size=10, color=COL_TEXT_CAPTION)


def slide_06_risks(prs, data, assets_dir):
    """Slide 6 — Riscos Ativos: top-N por severidade (crit + high), em 2 colunas.

    Design fiel ao artboard Paper `07 — Risks`:
      - Accent bar vertical colorida à esquerda (pela severidade)
      - Header com risk code + título
      - Tag "PROB · IMPACTO" à direita (com bg colorido por severidade)
      - Contramedida abaixo em cor de texto secundária
    Até 6 cards em 2 colunas (3x2).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, prs.slide_width, prs.slide_height, COL_BG_LIGHT)

    # Header
    add_text(slide, 80, 56, 400, 16, "06 · RISCOS",
             size=12, bold=True, color=COL_PRIMARY, tracking=0.2)
    add_logo(slide, 1180, 48, 60, assets_dir, "dark")
    add_text(slide, 80, 80, 900, 40,
             data.get("status", {}).get("risks_sentence", ""),
             size=28, color=COL_PRIMARY)
    add_text(slide, 80, 120, 1120, 20,
             "Riscos ativos com contramedidas em execução. Priorizados por severidade.",
             size=12, color=COL_TEXT_SUBTLE)

    # Filter risks: exclude upsides (is_upside=True) + keep only crit/high by default
    all_risks = data.get("risks", [])
    risks = [r for r in all_risks if not r.get("is_upside")
             and r.get("severity_class") in ("crit", "high")]
    # If fewer than 6 crit/high, include medium to fill the slide
    if len(risks) < 4:
        risks += [r for r in all_risks if not r.get("is_upside")
                  and r.get("severity_class") == "med"][: 6 - len(risks)]
    risks = risks[:6]

    if not risks:
        add_text(slide, 80, 340, 1120, 40,
                 "Nenhum risco crítico ou alto mapeado.",
                 size=18, color=COL_TEXT_SUBTLE, align="center")
    else:
        # Layout: 2 columns × up to 3 rows, card_w ~540 × card_h ~150
        card_w = 550
        card_h = 148
        gap_x = 20
        gap_y = 10
        start_x = 80
        start_y = 160
        for i, r in enumerate(risks):
            col = i % 2
            row = i // 2
            cx = start_x + col * (card_w + gap_x)
            cy = start_y + row * (card_h + gap_y)

            severity_label, tag_bg, bar_color = risk_severity_label(r)

            # Accent bar (6px wide, full height)
            add_rect(slide, cx, cy, 6, card_h, bar_color)
            # Card box
            add_rect(slide, cx + 6, cy, card_w - 6, card_h, COL_BG_LIGHT, line_rgb=COL_ROW_BORDER)

            # Risk code + title on first line
            code = r.get("code", "R")
            title = r.get("title", "")
            add_text(slide, cx + 20, cy + 12, card_w - 200, 30,
                     f"{code} — {truncate(title, 70)}",
                     size=13, bold=True, color=COL_PRIMARY)

            # Severity tag (top-right of card)
            tag_w = 130
            tag_x = cx + card_w - tag_w - 10
            add_rect(slide, tag_x, cy + 14, tag_w, 20, tag_bg)
            add_text(slide, tag_x + 4, cy + 16, tag_w - 8, 16,
                     severity_label, size=8, bold=True, color=COL_PRIMARY,
                     tracking=0.14, align="center")

            # Description (if present, up to 2 lines)
            desc = r.get("description", "")
            if desc:
                add_text(slide, cx + 20, cy + 48, card_w - 30, 32,
                         truncate(desc, 140),
                         size=10, color=COL_TEXT_SUBTLE)

            # Mitigation (bottom section)
            mitig = r.get("mitigation", "") or "— sem contramedida —"
            add_text(slide, cx + 20, cy + 92, 110, 14,
                     "CONTRAMEDIDA", size=8, bold=True,
                     color=COL_TEXT_SUBTLE, tracking=0.14)
            add_text(slide, cx + 20, cy + 108, card_w - 30, 36,
                     truncate(mitig, 150),
                     size=10, color=COL_TEXT_MUTED)

    # Footer: count summary
    hidden = len([r for r in all_risks if not r.get("is_upside")]) - len(risks)
    suffix = f" · +{hidden} com severidade menor não exibidos" if hidden > 0 else ""
    add_text(slide, 80, 670, 1120, 14,
             f"Fonte: 1-planning/artefatos/riscos.html · atualizado em {data.get('report_date','')}{suffix}",
             size=10, color=COL_TEXT_CAPTION)


def truncate(s, max_len):
    s = (s or "").strip()
    return s if len(s) <= max_len else s[: max_len - 1].rstrip() + "…"


def slide_07_closing(prs, data, assets_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, prs.slide_width, prs.slide_height, COL_PRIMARY)
    add_text(slide, 80, 64, 400, 20, "PRÓXIMOS PASSOS",
             size=13, bold=True, color=COL_ACCENT_LIME, tracking=0.2)
    add_logo(slide, 1180, 50, 60, assets_dir, "offwhite")

    ns = data.get("next_steps", [])
    primary = ns[0] if ns else None
    # Eyebrow
    add_text(slide, 80, 200, 600, 16, "UMA AÇÃO PRIORITÁRIA",
             size=11, bold=True, color=COL_BG_LIGHT, tracking=0.18)
    # Hero
    if primary:
        hero = f"{primary['action']} até {primary.get('deadline','')}"
    else:
        hero = "Sem ações cadastradas"
    add_text(slide, 80, 260, 1120, 120, hero,
             size=48, color=COL_BG_LIGHT)
    # Rationale
    if primary and primary.get("rationale"):
        add_text(slide, 80, 420, 900, 60, primary["rationale"],
                 size=14, color=COL_BG_LIGHT)
    # Divider
    add_rect(slide, 80, 600, 1120, 1, COL_BG_LIGHT)
    # Footer left
    add_text(slide, 80, 620, 600, 18,
             data.get("project", {}).get("footer_label", ""),
             size=12, color=COL_BG_LIGHT)
    # Footer right (contact)
    email = data.get("project", {}).get("pm_email") or "—"
    add_text(slide, 820, 620, 380, 18,
             f"Dúvidas: {email}", size=12, color=COL_BG_LIGHT, align="right")


# ---- Main ----

def build(data: dict, out_path: Path, assets_dir: Path, ctx: dict):
    """Builds the 7-slide status presentation. Ordering (v1.5):
    Cover · Agenda · Visão Geral · Roadmap · Mapa de Status Executivo · Riscos · Closing.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01_cover(prs, data, assets_dir)
    slide_02_agenda(prs, data, assets_dir)
    slide_03_visao_geral_roadmap(prs, data, assets_dir, ctx)
    slide_04_roadmap(prs, data, assets_dir, ctx)
    slide_05_mapa_status_executivo(prs, data, assets_dir, ctx)
    slide_06_risks(prs, data, assets_dir)
    slide_07_closing(prs, data, assets_dir)

    prs.save(str(out_path))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True, type=Path)
    ap.add_argument("--assets-dir", type=Path)
    ap.add_argument("--out-dir", required=True, type=Path)
    ap.add_argument(
        "--roadmap-html", type=Path,
        help="Path to 1-planning/artefatos/roadmap-marcos.html (source for swim-lane screenshots).",
    )
    args = ap.parse_args()

    script_dir = Path(__file__).parent
    skill_dir = script_dir.parent
    assets_dir = args.assets_dir or (skill_dir / "assets")

    if not args.data.exists():
        print(f"✗ Arquivo de dados não encontrado: {args.data}", file=sys.stderr)
        sys.exit(2)
    args.out_dir.mkdir(parents=True, exist_ok=True)

    data = json.loads(args.data.read_text(encoding="utf-8"))

    # Context passed through to slide builders (tmp dir for screenshots, HTML sources, dates)
    tmp_dir = Path(tempfile.mkdtemp(prefix="status-materials-"))
    report_date = None
    try:
        report_date = datetime.strptime(data.get("report_date", ""), "%Y-%m-%d")
    except ValueError:
        pass
    ctx = {
        "tmp_dir": tmp_dir,
        "roadmap_html": args.roadmap_html,
        "report_date": report_date,
    }

    out_path = args.out_dir / "status-presentation.pptx"
    build(data, out_path, assets_dir, ctx)
    print(f"✓ PPTX gerado: {out_path}")

    if data.get("warnings"):
        print(f"⚠ {len(data['warnings'])} warning(s):", file=sys.stderr)
        for w in data["warnings"]:
            print(f"   · {w}", file=sys.stderr)


if __name__ == "__main__":
    main()
